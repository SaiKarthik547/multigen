"""
PresentationEngine — PowerPoint generation with style-theme support.

Phase 1: Migrated from MultiGenAi.py with ExecutionContext injection.
Phase 7: Will add LLM-planned slide structure, visual embedding reuse,
         auto-image insertion, and layout consistency.
"""

from __future__ import annotations

import pathlib
import re
from dataclasses import dataclass
from typing import TYPE_CHECKING, Optional

from multigenai.core.logging.logger import get_logger

if TYPE_CHECKING:
    from multigenai.core.execution_context import ExecutionContext
    from multigenai.llm.schema_validator import DocumentGenerationRequest

LOG = get_logger(__name__)


@dataclass
class PresentationResult:
    """Output from the PresentationEngine."""
    path: str
    title: str
    slide_count: int
    success: bool = True
    error: Optional[str] = None


class PresentationEngine:
    """
    Generates PowerPoint presentations from a natural language prompt.

    Usage:
        engine = PresentationEngine(ctx)
        result = engine.run(DocumentGenerationRequest(prompt="quantum computing", output_format="pptx"))
    """

    def __init__(self, ctx: "ExecutionContext") -> None:
        self._ctx = ctx
        self._out_dir = pathlib.Path(ctx.settings.output_dir)
        self._out_dir.mkdir(parents=True, exist_ok=True)

    def run(self, request: "DocumentGenerationRequest") -> PresentationResult:
        import hashlib
        slug = re.sub(r"[^A-Za-z0-9\-_]+", "_", request.prompt)[:40]
        slug += f"_{hashlib.sha1(request.prompt.encode()).hexdigest()[:8]}"
        out_path = self._out_dir / f"{slug}.pptx"

        title, content = self._fetch_content(request.prompt)

        try:
            import pptx as python_pptx
            from pptx.util import Pt
        except ImportError:
            msg = "python-pptx not installed — run: pip install python-pptx"
            LOG.error(msg)
            return PresentationResult(path="", title=title, slide_count=0, success=False, error=msg)

        prs = python_pptx.Presentation()

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = f"Presentation about: {request.prompt}"

        # Content slides (4 sentences per slide)
        sentences = re.split(r"(?<=[.!?])\s+", content)
        per_slide = 4
        slide_count = 1
        for i in range(0, len(sentences), per_slide):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Key Points ({i // per_slide + 1})"
            tf = slide.shapes.placeholders[1].text_frame
            tf.clear()
            for sentence in sentences[i:i + per_slide]:
                p = tf.add_paragraph()
                p.text = sentence
                p.level = 1
            slide_count += 1

        prs.save(str(out_path))
        LOG.info(f"Presentation saved: {out_path} ({slide_count} slides)")
        return PresentationResult(path=str(out_path), title=title, slide_count=slide_count)

    def _fetch_content(self, query: str) -> tuple[str, str]:
        try:
            import wikipediaapi
            wiki = wikipediaapi.Wikipedia("MultiGenAI/1.0", "en")
            page = wiki.page(query)
            if page.exists():
                return page.title, " ".join(page.summary.split()[:400]) + "…"
        except Exception as exc:
            LOG.warning(f"Wikipedia fetch: {exc}")
        return f"Presentation: {query}", f"A comprehensive look at {query}."

    def run_with_llm_planning(self, request: "DocumentGenerationRequest") -> PresentationResult:
        """[Phase 7] LLM-driven slide structure with cinematic theme."""
        raise NotImplementedError("LLM presentation planning activates in Phase 7.")

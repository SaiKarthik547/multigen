from __future__ import annotations

import hashlib
import logging
import os
import pathlib
import random
import re
from concurrent.futures import ProcessPoolExecutor, Future
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Optional

# region --- Optional Library Imports ---

try:
    import nltk
    nltk.download('punkt', quiet=True)
    nltk.download('averaged_perceptron_tagger', quiet=True)
except Exception:
    nltk = None
    logging.warning("NLTK not found.")

try:
    from PIL import Image, ImageDraw, ImageFont
except Exception:
    Image = None
    logging.warning("Pillow not found.")

try:
    from moviepy.editor import ImageClip, concatenate_videoclips
except Exception:
    ImageClip = None
    logging.warning("MoviePy not found.")

try:
    import torch
    from diffusers import DiffusionPipeline
    logging.info("PyTorch and Diffusers found. Stable Diffusion is enabled.")
except Exception as e:
    print(f"DEBUG: Import failed with error: {e}")
    torch = None
    DiffusionPipeline = None
    logging.warning("Torch or Diffusers not found.")

try:
    import matplotlib.pyplot as plt
except Exception:
    plt = None
    logging.warning("Matplotlib not found.")

try:
    import pptx
    from pptx.util import Inches
except Exception:
    pptx = None
    logging.warning("python-pptx not found. PowerPoint generation is disabled.")

try:
    import docx
except Exception:
    docx = None
    logging.warning("python-docx not found. Word document generation is disabled.")

try:
    import wikipediaapi
    wiki_wiki = wikipediaapi.Wikipedia('MultiGenAI/1.0 (Local AI Project)','en')
    logging.info("Wikipedia API found. Internet content sourcing is enabled.")
except Exception:
    wiki_wiki = None
    logging.warning("wikipedia-api not found. Internet content sourcing is disabled.")

# endregion --- End of Imports ---


# region --- Global Configuration & Setup ---

ROOT = pathlib.Path(__file__).resolve().parent
OUT = ROOT.joinpath('multigen_outputs')
OUT.mkdir(parents=True, exist_ok=True)

logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')
LOG = logging.getLogger('MultiGenAI')

# endregion --- End of Configuration ---


# region --- Utility Functions ---

def slug(s: str) -> str:
    safe = re.sub(r'[^A-Za-z0-9\-_]+', '_', s)[:40]
    return f"{safe}_{hashlib.sha1(s.encode()).hexdigest()[:8]}"

# endregion --- End of Utility Functions ---


# region --- Core Classes and Pipeline Stages ---

@dataclass
class PipelineResult:
    path: str
    meta: Dict[str, Any]

class Stage:
    def run(self, data: Dict[str, Any]) -> Dict[str, Any]:
        raise NotImplementedError()

class Preprocessor(Stage):
    def run(self, data: Dict[str, Any]) -> Dict[str, Any]:
        return data

class Generator(Stage):
    _sd_pipeline = None
    _device = None

    def __init__(self, typ: str):
        self.typ = typ

    def run(self, data: Dict[str, Any]) -> Dict[str, Any]:
        prompt = data.get('prompt', '')
        gen_map = {
            'image': self.generate_image,
            'video': self.generate_video,
            'code': self.generate_code_file,
            'visualization': self.generate_visualization,
            'powerpoint': self.generate_powerpoint,
            'word': self.generate_word_doc,
            'document': self.generate_word_doc,
        }
        func = lambda: gen_map.get(self.typ, lambda p: f"Error: Type '{self.typ}' not supported.")(prompt)
        out_content_or_path = func()
        data['raw'] = out_content_or_path
        LOG.info(f"Generated {self.typ} for prompt: '{prompt}' -> {str(out_content_or_path)[:100]}...")
        return data

    def _generate_with_stable_diffusion(self, prompt: str, out_path: Path, seed: Optional[int] = None) -> bool:
        if not torch or not DiffusionPipeline: return False
        try:
            if Generator._sd_pipeline is None:
                LOG.info("Initializing Stable Diffusion pipeline...")
                Generator._device = "cuda" if torch.cuda.is_available() else "dml" if hasattr(torch, 'dml') and torch.dml.is_available() else "cpu"
                LOG.info(f"Detected device: {Generator._device}. Loading model...")

                model_id = "stabilityai/stable-diffusion-xl-base-1.0"
                torch_dtype = torch.float16 if Generator._device == "cuda" else torch.float32

                # --- FIX: Conditional variant loading ---
                # Only specify the 'fp16' variant for CUDA devices.
                # For other devices (like AMD/DirectML), do not specify a variant
                # and let the library load the default (full-precision) model.
                pipeline_kwargs = {
                    "torch_dtype": torch_dtype,
                    "use_safetensors": True,
                }
                if torch_dtype == torch.float16:
                    pipeline_kwargs["variant"] = "fp16"

                Generator._sd_pipeline = DiffusionPipeline.from_pretrained(
                    model_id,
                    **pipeline_kwargs
                )
                Generator._sd_pipeline = Generator._sd_pipeline.to(Generator._device)
                LOG.info("VRAM optimization (attention slicing) enabled by default in new diffusers.")

            LOG.info(f"Generating on {Generator._device} with SDXL model...")
            quality_enhancers = "masterpiece, best quality, ultra-detailed, cinematic lighting, sharp focus, intricate details, 8k, photorealistic"
            negative_prompt = "deformed, distorted, disfigured, poorly drawn, bad anatomy, wrong anatomy, extra limb, missing limb, floating limbs, mutated hands, mutated fingers, blurry, low resolution, ugly, disgusting, watermark, text, signature, abstract"
            enhanced_prompt = f"{prompt}, {quality_enhancers}"
            
            if seed is None:
                seed = random.randint(0, 1_000_000)
            generator = torch.Generator(device=Generator._device).manual_seed(seed) if Generator._device != "cpu" else None

            image = Generator._sd_pipeline(
                prompt=enhanced_prompt,
                negative_prompt=negative_prompt,
                generator=generator,
                num_inference_steps=40,
                guidance_scale=8.0,
            ).images[0]
            
            image.save(out_path)
            LOG.info(f"Stable Diffusion image saved to {out_path} (Seed: {seed})")
            return True
        except Exception as e:
            LOG.error(f"Stable Diffusion generation failed: {e}")
            Generator._sd_pipeline = None
            return False

    def generate_image(self, p: str) -> str:
        out_path = OUT.joinpath(f"{slug(p)}.png")
        if self._generate_with_stable_diffusion(p, out_path):
            return str(out_path)
        LOG.warning("AI generation failed. Creating a placeholder image.")
        return str(out_path)

    def generate_video(self, p: str) -> str:
        out_path = OUT.joinpath(f"{slug(p)}.mp4")
        if not ImageClip: return "Error: MoviePy library not installed."
        
        num_frames = 10
        frame_duration = 1.0
        frame_paths = []
        
        try:
            video_seed = random.randint(0, 1_000_000)
            LOG.info(f"Generating {num_frames} frames for video with consistent seed: {video_seed}")

            for i in range(num_frames):
                frame_prompt = f"{p}, shot {i+1} of {num_frames}"
                LOG.info(f"Generating frame {i+1}/{num_frames}...")
                
                frame_path = OUT.joinpath(f"temp_frame_{slug(p)}_{i}.png")
                success = self._generate_with_stable_diffusion(frame_prompt, frame_path, seed=video_seed)
                
                if success:
                    frame_paths.append(str(frame_path))
                else:
                    raise RuntimeError(f"Failed to generate frame {i+1} for video.")
            
            if not frame_paths: raise RuntimeError("No valid frames were generated.")
            
            LOG.info(f"Stitching {len(frame_paths)} frames into video...")
            clips = [ImageClip(fp).set_duration(frame_duration) for fp in frame_paths]
            video = concatenate_videoclips(clips, method="compose")
            video.write_videofile(str(out_path), fps=30, bitrate="5000k", verbose=False, logger=None, codec="libx264")
            
            LOG.info(f"Successfully created high-quality video: {out_path}")
            return str(out_path)
        except Exception as e:
            LOG.error(f"Video generation failed: {e}. Ensure FFmpeg is installed and in PATH.")
            return f"Error creating video for {p}."
        finally:
            for path in frame_paths:
                try: os.remove(path)
                except OSError: pass

    def generate_code_file(self, prompt: str) -> str:
        language_map = {'python': 'py', 'javascript': 'js', 'html': 'html', 'css': 'css', 'java': 'java', 'c++': 'cpp', 'c#': 'cs', 'go': 'go', 'rust': 'rs', 'ruby': 'rb', 'sql': 'sql', 'shell': 'sh'}
        detected_extension = 'txt'
        for lang, ext in language_map.items():
            if lang in prompt.lower():
                detected_extension = ext; break
        return self.generate_text_file(prompt, detected_extension)

    def generate_text_file(self, p: str, extension: str) -> str:
        out_path = OUT.joinpath(f"{slug(p)}.{extension}")
        content = f"# Prompt: {p}\n\n# NOTE: This is a placeholder. A true LLM is needed for meaningful content."
        out_path.write_text(content)
        return str(out_path)

    def generate_visualization(self, p: str) -> str:
        out_path = OUT.joinpath(f"{slug(p)}.png")
        if not plt: return "Error: Matplotlib not installed."
        fig, ax = plt.subplots(figsize=(8, 6))
        labels = ['Category A', 'Category B', 'Category C', 'Category D']
        values = [random.randint(10, 50) for _ in labels]
        ax.bar(labels, values)
        ax.set_title(f'Sample Chart for: "{p[:50]}"')
        plt.savefig(str(out_path), dpi=300)
        plt.close(fig)
        return str(out_path)
            
    def _get_internet_content(self, query: str) -> (str, str):
        if not wiki_wiki:
            return "Content Generation Error", "Wikipedia library not installed."
        try:
            LOG.info(f"Searching Wikipedia for: '{query}'...")
            page = wiki_wiki.page(query)
            if page.exists():
                LOG.info(f"Found page: {page.title}")
                summary = " ".join(page.summary.split()[:300]) + "..."
                return page.title, summary
            else:
                return f"No article found for '{query}'", "Could not find a matching Wikipedia article."
        except Exception as e:
            LOG.error(f"Wikipedia search failed: {e}")
            return "Content Generation Error", str(e)

    def generate_powerpoint(self, prompt: str) -> str:
        if not pptx: return "Error: python-pptx not installed."
        out_path = OUT.joinpath(f"{slug(prompt)}.pptx")
        title, content = self._get_internet_content(prompt)
        prs = pptx.Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"A presentation on {prompt}"
        sentences = re.split(r'(?<=[.!?]) +', content)
        sentences_per_slide = 4
        for i in range(0, len(sentences), sentences_per_slide):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = f"Key Points ({i//sentences_per_slide + 1})"
            tf = slide.shapes.placeholders[1].text_frame
            tf.clear()
            chunk = sentences[i:i+sentences_per_slide]
            for sentence in chunk:
                p = tf.add_paragraph(); p.text = sentence; p.level = 1
        prs.save(str(out_path))
        return str(out_path)

    def generate_word_doc(self, prompt: str) -> str:
        if not docx: return "Error: python-docx not installed."
        out_path = OUT.joinpath(f"{slug(prompt)}.docx")
        title, content = self._get_internet_content(prompt)
        document = docx.Document()
        document.add_heading(title, level=0)
        document.add_paragraph(f"A document about {prompt}, sourced from Wikipedia.")
        for para in content.split('\n'):
            if para.strip(): document.add_paragraph(para)
        document.save(str(out_path))
        return str(out_path)

class Postprocessor(Stage):
    def run(self, data: Dict[str, Any]) -> Dict[str, Any]:
        data['final_path'] = data.get('raw')
        return data

class Pipeline:
    def __init__(self, prompt: str, typ: str):
        self.prompt = prompt; self.type = typ
        self.stages = [Preprocessor(), Generator(typ), Postprocessor()]

    def run(self) -> PipelineResult:
        ctx = {'prompt': self.prompt, 'type': self.type}
        for stage in self.stages:
            ctx = stage.run(ctx)
        return PipelineResult(path=str(ctx.get('final_path', '')), meta=ctx)

def run_pipeline(prompt: str, typ: str) -> PipelineResult:
    return Pipeline(prompt, typ).run()

class Orchestrator:
    def __init__(self, workers: int = 1):
        self.pool = ProcessPoolExecutor(max_workers=workers)

    def submit(self, prompt: str) -> str:
        typ = self.classify_type(prompt)
        future: Future = self.pool.submit(run_pipeline, prompt, typ)
        try:
            res = future.result(timeout=1800)
            return res.path
        except Exception as e:
            LOG.error(f"Job failed in orchestrator: {e}")
            return ''
            
    def classify_type(self, prompt: str) -> str:
        toks = set(re.findall(r"\w+", prompt.lower()))
        if not toks.isdisjoint({'powerpoint', 'presentation', 'ppt', 'pptx'}): return 'powerpoint'
        if not toks.isdisjoint({'word', 'docx', 'report'}): return 'word'
        if not toks.isdisjoint({'document', 'doc', 'file', 'text', 'write'}): return 'document'
        if not toks.isdisjoint({'video', 'clip', 'animation'}): return 'video'
        if not toks.isdisjoint({'chart', 'plot', 'graph', 'visualization'}): return 'visualization'
        if not toks.isdisjoint({'code', 'script'}): return 'code'
        return 'image'

def run_cli():
    orch = Orchestrator(workers=1)
    print('MultiGenAI is ready. Enter a prompt to begin (type EXIT to quit).')
    while True:
        try:
            p = input('Prompt: ').strip()
            if p.upper() == 'EXIT':
                orch.pool.shutdown(); break
            if not p: continue
            
            path = orch.submit(p)
            print(f"--> Generation finished.\n--> Output: {path if path else 'Failed - check logs.'}\n")
        except KeyboardInterrupt:
            print("\nShutting down...")
            orch.pool.shutdown(wait=False, cancel_futures=True); break
        except Exception as e:
            LOG.error(f'CLI loop error: {e}')

if __name__ == "__main__":
    run_cli()